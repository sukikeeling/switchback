# -*- coding: utf-8 -*-
"""Generate the GOAI Track-1 submission deck (16:9) with python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RAIL = RGBColor(0x12, 0x3A, 0x5F)      # 深轨蓝
RAIL2 = RGBColor(0x0C, 0x2C, 0x49)
GREEN = RGBColor(0x2F, 0x7D, 0x5A)     # 验证绿
COPPER = RGBColor(0xB0, 0x76, 0x3C)    # 历史铜
INK = RGBColor(0x1C, 0x27, 0x33)
MUTED = RGBColor(0x5A, 0x6B, 0x7A)
BG = RGBColor(0xF6, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

FONT = "Microsoft YaHei"


def box(slide, x, y, w, h, fill=WHITE, line=RAIL):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def text(slide, x, y, w, h, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = runs if isinstance(runs, list) else [(runs, {})]
    for i, (txt, style) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", align)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(style.get("size", size))
        r.font.bold = style.get("bold", bold)
        r.font.color.rgb = style.get("color", color)
        r.font.name = style.get("font", font)
    return tb


def header(slide, num, title):
    box(slide, 0, 0, 13.333, 0.9, fill=RAIL)
    text(slide, 0.5, 0.12, 12, 0.7, [
        (f"{num}  ", {"color": COPPER, "bold": True, "size": 18}),
        (title, {"color": WHITE, "bold": True, "size": 18}),
    ])
    text(slide, 11.4, 0.25, 1.6, 0.5, "Switchback", size=11, color=WHITE, align=PP_ALIGN.RIGHT)


# ---------------- S1 封面 ----------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=RAIL2)
box(s, 0, 7.28, 13.333, 0.22, fill=COPPER)
text(s, 1, 1.1, 11, 0.5, "GOAI WORLD AI OPEN SOURCE COMPETITION 2026 · HANGZHOU · TRACK 1 · AGENT INFRA", size=12, color=RGBColor(0x7C, 0xE7, 0xA8), bold=True)
text(s, 1, 2.0, 11, 1.2, "折返治理 · Switchback", size=52, color=WHITE, bold=True)
text(s, 1, 3.2, 11, 0.6, "Switchback Governance — the Human-in-the-Loop Governance Layer for Multi-Agent Teams", size=16, color=RGBColor(0xB9, 0xC6, 0xD2))
text(s, 1, 4.2, 11, 0.6, "让每一次 Agent 决策，都沿轨道可查、可停、可回头。", size=20, color=COPPER)
text(s, 1, 5.3, 11, 0.5, "赛道① 新智基座 Agent Infra ｜ 参赛 sukikeeling ｜ 仓库 github.com/sukikeeling/switchback", size=13, color=WHITE)
text(s, 1, 5.9, 11, 0.5, "旗舰证据：百年京张 AI 创新带开源征集 84/100 最高纪录（open-city-ai/haidian）", size=13, color=RGBColor(0xB9, 0xC6, 0xD2))

# ---------------- S2 问题与场景 ----------------
s = prs.slides.add_slide(BLANK)
header(s, "一", "问题与场景：Agent 从 Demo 到 Production 的四道坎")
rows = [
    ("结果不可信", "数值三处对不齐、引用错位、模型幻觉", "证据核验 Skill：三处对齐 + 双向引用 + 内容哈希"),
    ("过程不可停", "编排一旦开始就自动往下走，错误放大", "折返点：到站必停，复核后才继续"),
    ("出事不可回", "没有审批/回滚/隔离语义，只能硬杀进程", "道岔三态：正线/侧线/入段，不设自动恢复"),
    ("账目不可审", "谁做了什么、基于哪个数据版本，无账本", "K 标账本：内容寻址哈希链，逐条可复核"),
]
y = 1.3
for name, symptom, fix in rows:
    box(s, 0.6, y, 12.1, 1.0, fill=WHITE, line=GREEN)
    text(s, 0.9, y + 0.2, 2.2, 0.6, name, size=15, bold=True, color=RAIL)
    text(s, 3.2, y + 0.2, 4.4, 0.6, symptom, size=12, color=INK)
    text(s, 7.9, y + 0.2, 4.5, 0.6, fix, size=12, color=GREEN, bold=True)
    y += 1.2
box(s, 0.6, 6.3, 12.1, 0.9, fill=RGBColor(0xEE, 0xF5, 0xF0), line=GREEN)
text(s, 0.9, 6.45, 11.5, 0.6, "真实场景：京张开源征集 16 个 PR 亲历全部事故 → 84/100 最高纪录", size=13, bold=True, color=GREEN)

# ---------------- S3 核心洞察 ----------------
s = prs.slides.add_slide(BLANK)
header(s, "二", "核心洞察：铁路的『人字坡』智慧")
box(s, 0.6, 1.3, 12.1, 1.7, fill=RGBColor(0xFA, 0xF1, 0xE6), line=COPPER)
text(s, 0.9, 1.45, 11.5, 0.4, "制度原型：青龙桥『人字形折返』展线（1908 · 詹天佑）", size=14, bold=True, color=COPPER)
text(s, 0.9, 2.0, 11.5, 0.8, "列车无法直爬 33‰ 陡坡，必须在折返点停车、换向、以退为进。运行 118 年零重大事故——在能力极限前主动停下重估，比硬闯更接近最优解。", size=13, color=INK)
text(s, 0.6, 3.3, 12.1, 1.2, [
    ("核心洞察：", {"bold": True, "size": 16, "color": RAIL}),
    ("多 Agent 协同系统与陡坡列车同构——任务爬坡时，Agent 的自回归式自动续行就是『硬闯陡坡』。治理的关键不是更强的模型，而是制度化的折返点：让系统『到站必停、复核、可回头』。", {"size": 15}),
], color=INK)
mech = [
    ("折返点", "到站必停，任何否决即强制折返"),
    ("坡度分级", "缓/中/陡，越陡复核越严"),
    ("K 标版本", "不可变内容寻址哈希链"),
    ("道岔三态", "正线/侧线/入段，不设自动恢复"),
]
x = 0.6
for name, desc in mech:
    box(s, x, 4.9, 2.95, 1.5, fill=WHITE, line=RAIL)
    text(s, x + 0.2, 5.1, 2.6, 0.5, name, size=15, bold=True, color=RAIL)
    text(s, x + 0.2, 5.7, 2.6, 0.6, desc, size=11, color=INK)
    x += 3.05

# ---------------- S4 方案设计 ----------------
s = prs.slides.add_slide(BLANK)
header(s, "三", "方案设计：折返治理协议")
text(s, 0.6, 1.2, 12, 0.4, "架构：治理层（G-GOVERNOR 监理）→ AgentTeams 编排层（Manager-Worker）→ 基础设施（状态/记忆/可观测/MCP）", size=14, bold=True, color=RAIL)
box(s, 0.6, 1.9, 12.1, 0.9, fill=RGBColor(0xEE, 0xF5, 0xF0), line=GREEN)
text(s, 0.9, 2.05, 11.5, 0.6, "G-GOVERNOR 监理 Agent：折返点裁决 · 坡度准入 · 道岔三态 · K 标账本 · 不设自动恢复", size=14, bold=True, color=GREEN)
agents = [("M-OPER", "运营总监 Manager"), ("W-RESEARCH", "调研/来源"), ("W-CREATE", "内容/双语"), ("W-VERIFY", "校验/复算"), ("W-RISK", "风险/合规")]
x = 0.6
for a, r in agents:
    box(s, x, 3.0, 2.4, 1.2, fill=WHITE, line=RAIL)
    text(s, x + 0.15, 3.2, 2.1, 0.5, a, size=14, bold=True, color=RAIL)
    text(s, x + 0.15, 3.7, 2.1, 0.5, r, size=11, color=INK)
    x += 2.45
box(s, 0.6, 4.4, 12.1, 2.4, fill=WHITE, line=COPPER)
text(s, 0.9, 4.55, 11.5, 0.4, "协议不变量（可复算、可复核、可接手）", size=14, bold=True, color=COPPER)
inv = [
    "① 到站必停：折返点未裁决，任务不得继续执行",
    "② 任何否决即折返：一票 TURN_BACK/DEPOT ⇒ 强制折返，绝不放行",
    "③ 坡度越高越严：陡坡需三方（责任人/专业/公众）复核",
    "④ K 标不可篡改：SHA-256 哈希链，篡改即断链",
    "⑤ 不设自动恢复：被放行又入段的任务禁止自动回正线",
]
for i, invt in enumerate(inv):
    text(s, 1.0, 5.0 + i * 0.34, 11.3, 0.3, invt, size=12, color=INK)

# ---------------- S5 闭环与 Skill ----------------
s = prs.slides.add_slide(BLANK)
header(s, "四", "闭环 8 步 + 六大 Skill")
steps = [("①任务输入", "M-OPER 接单"), ("②任务拆解", "坡度分级 DAG"), ("③上下文传递", "Matrix+SharedState"), ("④工具调用", "Higress+MCP"),
         ("⑤结果验证", "W-VERIFY"), ("⑥证据沉淀", "K 标 seal"), ("⑦审批回滚", "折返点三方复核"), ("⑧经验沉淀", "lessons-learned")]
x = 0.6
for st, who in steps:
    box(s, x, 1.3, 1.5, 1.4, fill=WHITE, line=GREEN)
    text(s, x + 0.05, 1.45, 1.4, 0.7, st, size=10.5, bold=True, color=RAIL)
    text(s, x + 0.05, 2.15, 1.4, 0.5, who, size=9, color=MUTED)
    x += 1.55
text(s, 0.6, 3.0, 12, 0.4, "六大可执行 Skill（switchback/skills.py，每个带机器可读 SkillSpec）", size=14, bold=True, color=RAIL)
skills = [("grade-access", "坡度准入"), ("evidence-verify", "证据核验"), ("kmarker-ledger", "K 标账本"),
          ("risk-ledger", "风险台账"), ("switch-decision", "折返裁决"), ("lessons-learned", "经验沉淀")]
grid = [skills[0:3], skills[3:6]]
yy = 3.5
for row in grid:
    xx = 0.6
    for name, desc in row:
        box(s, xx, yy, 4.0, 0.9, fill=WHITE, line=RAIL)
        text(s, xx + 0.15, yy + 0.1, 2.2, 0.4, name, size=13, bold=True, color=RAIL)
        text(s, xx + 2.3, yy + 0.15, 1.6, 0.6, desc, size=11, color=INK)
        xx += 4.1
    yy += 1.1
box(s, 0.6, 5.9, 12.1, 1.0, fill=RGBColor(0xEE, 0xF5, 0xF0), line=GREEN)
text(s, 0.9, 6.05, 11.5, 0.7, "以 AgentTeams（Hiclaw）为设计基点：Manager 控制流 + Worker 任务流 + Matrix 房间人机可见可干预 + 凭证透传；Agent Identity 清单见 docs/identity.md", size=13, bold=True, color=GREEN)

# ---------------- S6 真实案例 ----------------
s = prs.slides.add_slide(BLANK)
header(s, "五", "真实案例：京张 84 分（可复现）")
rows = [("v5", "PR#605", "67", "PASS"), ("v8", "PR#1220", "70", "折返↩"), ("v8.1", "PR#1468", "84", "PASS"), ("v8.2", "PR#1816", "70", "折返↩"), ("v8.5", "PR#2205", "77", "折返↩"), ("v8.10", "PR#2328", "76", "折返↩")]
y = 1.4
text(s, 0.6, 1.1, 12, 0.4, "python -m switchback.cli replay jingzhang", size=13, color=RAIL, font="Consolas")
text(s, 0.6, y, 1.5, 0.4, "版本", size=12, bold=True, color=WHITE)
text(s, 2.3, y, 1.5, 0.4, "PR", size=12, bold=True, color=WHITE)
text(s, 4.0, y, 1.0, 0.4, "分", size=12, bold=True, color=WHITE)
text(s, 5.2, y, 2.0, 0.4, "裁决", size=12, bold=True, color=WHITE)
box(s, 0.6, y, 12.1, 0.42, fill=RAIL)
y += 0.5
for v, pr, sc, verdict in rows:
    box(s, 0.6, y, 12.1, 0.42, fill=WHITE if y % 2 == 0 else BG, line=GREEN)
    color = GREEN if verdict == "PASS" else COPPER
    text(s, 0.6, y + 0.05, 1.5, 0.3, v, size=11, bold=True, color=RAIL)
    text(s, 2.3, y + 0.05, 1.5, 0.3, pr, size=11, color=INK)
    text(s, 4.0, y + 0.05, 1.0, 0.3, sc, size=11, bold=True, color=INK)
    text(s, 5.2, y + 0.05, 2.0, 0.3, verdict, size=11, bold=True, color=color)
    y += 0.5
box(s, 0.6, y + 0.2, 12.1, 1.0, fill=RGBColor(0xEE, 0xF5, 0xF0), line=GREEN)
text(s, 0.9, y + 0.35, 11.5, 0.7, "每次提交过『证据核验 → 折返点复核 → K 标放行/折返』，账本不可变、Trace 可回放。真实任务闭环 + 结果校验 + 安全熔断审计。", size=13, bold=True, color=GREEN)

# ---------------- S7 落地与开源 ----------------
s = prs.slides.add_slide(BLANK)
header(s, "六", "可行性与开放开源")
box(s, 0.6, 1.3, 12.1, 1.1, fill=WHITE, line=RAIL)
text(s, 0.9, 1.45, 11.5, 0.4, "工程落地", size=13, bold=True, color=RAIL)
text(s, 0.9, 1.85, 11.5, 0.5, "Python 3.10+ 零依赖可运行 · 43 项 pytest 全绿 · GitHub Actions CI（3.10/3.11/3.12）· CLI 一键复现", size=12, color=INK)
box(s, 0.6, 2.6, 12.1, 1.1, fill=WHITE, line=RAIL)
text(s, 0.9, 2.75, 11.5, 0.4, "落地路线", size=13, bold=True, color=RAIL)
text(s, 0.9, 3.15, 11.5, 0.5, "初赛（方案设计）→ 复赛 9.3（AgentTeams 集群部署 + Demo）→ 决赛 9.22（现场 Demo + 审计链路展示）", size=12, color=INK)
box(s, 0.6, 3.9, 12.1, 1.1, fill=WHITE, line=RAIL)
text(s, 0.9, 4.05, 11.5, 0.4, "开放/开源", size=13, bold=True, color=RAIL)
text(s, 0.9, 4.45, 11.5, 0.5, "Apache-2.0 全量开源 · 依赖边界清晰（运行零依赖）· 京张案例数据来自公开开源竞赛 · 权利逐资产登记", size=12, color=INK)
box(s, 0.6, 5.3, 12.1, 1.1, fill=WHITE, line=RAIL)
text(s, 0.9, 5.45, 11.5, 0.4, "评审指标对照", size=13, bold=True, color=RAIL)
text(s, 0.9, 5.85, 11.5, 0.5, "场景价值 25% ｜ 多 Agent 闭环 25% ｜ Skill 工程 25% ｜ 工程落地安全审计 20% ｜ 开源 5% —— 全部命中", size=12, color=INK)
box(s, 0.6, 6.6, 12.1, 0.7, fill=RAIL)
text(s, 0.9, 6.78, 11.5, 0.4, "让每一次 Agent 决策，都沿轨道可查、可停、可回头。", size=15, bold=True, color=WHITE)

prs.save("D:/switchback/competition/GOAI-初赛方案-折返治理Switchback.pptx")
print("PPTX saved")
